import pandas as pd
import pytesseract
import fitz
import pathlib
from pptx import Presentation
from docx import Document
from . import config

def parse_text_from_file(file_path: str) -> str:
    """Extracts text content from various file types."""
    suffix = pathlib.Path(file_path).suffix.lower()
    text_content = ""
    try:
        if suffix in ['.png', '.jpg', '.jpeg']:
            text_content = pytesseract.image_to_string(str(file_path))
        elif suffix == '.csv':
            df = pd.read_csv(file_path)
            text_content = df.to_string()
        elif suffix == '.xlsx':
            df_dict = pd.read_excel(file_path, sheet_name=None)
            for sheet_name, df in df_dict.items():
                text_content += f"--- Sheet: {sheet_name} ---\n{df.to_string()}\n\n"
        elif suffix == '.pdf':
            with fitz.open(file_path) as doc:
                for page in doc:
                    text_content += page.get_text()
        elif suffix == '.pptx':
            prs = Presentation(file_path)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text_content += shape.text + "\n"
        elif suffix == '.docx':
            doc = Document(file_path)
            for para in doc.paragraphs:
                text_content += para.text + "\n"
            for table in doc.tables:
                for row in table.rows:
                    for cell in row.cells:
                        text_content += cell.text + "\t"
                    text_content += "\n"
        else:
            return f"No parser available for file type: {suffix}"
        return text_content if text_content.strip() else "No text content found."
    except Exception as e:
        return f"Error parsing file {file_path}: {e}"