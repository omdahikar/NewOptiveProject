'''
import cv2
import fitz
import pandas as pd
import pytesseract
import re
import os
import spacy
from pptx import Presentation
from docx import Document
from . import config, audit_log

# Load AI model
try:
    NLP = spacy.load("en_core_web_sm")
except OSError:
    print("spaCy model 'en_core_web_sm' not found. Please run 'python -m spacy download en_core_web_sm'")
    NLP = None

pytesseract.pytesseract.tesseract_cmd = config.TESSERACT_CMD_PATH

# --- HELPER & CORE LOGIC ---

def _is_text_sensitive_advanced(text: str) -> bool:
    """A quick check to see if text likely contains any sensitive information."""
    if not text or not isinstance(text, str): return False
    # Check regex first as it's the fastest
    for pattern in config.PII_REDACTION_RULES.keys():
        if pattern.search(text): return True
    # Then check for specific keywords
    for keyword in config.SENSITIVE_KEYWORDS:
        if re.search(r'\b' + re.escape(keyword) + r'\b', text, re.IGNORECASE): return True
    # Finally, use the AI model for a deeper check
    if NLP:
        doc = NLP(text)
        for ent in doc.ents:
            if ent.label_ in config.NER_REDACTION_RULES: return True
    return False

def _apply_contextual_rules(text: str, file_name: str, nlp_doc) -> str:
    """Applies rules that redact information based on the text that precedes it."""
    redacted_text = text
    for rule in config.CONTEXTUAL_RULES:
        trigger_phrase = rule["trigger_phrase"]
        for match in re.finditer(re.escape(trigger_phrase), redacted_text, re.IGNORECASE):
            following_token_start = match.end()
            for token in nlp_doc:
                if token.idx >= following_token_start:
                    if token.pos_ == rule["TARGET_TYPE"] or (token.like_num and rule["TARGET_TYPE"] == "NUM"):
                        tag = rule["redaction_tag"]
                        start, end = token.idx, token.idx + len(token.text)
                        if "<" not in redacted_text[start:end]:
                            audit_log.log_redaction(file_name, token.text, tag, f"Contextual: '{trigger_phrase}'")
                            redacted_text = redacted_text[:start] + tag + redacted_text[end:]
                        break # Move to the next trigger phrase match
    return redacted_text

def redact_text_pipeline(text: str, file_name: str) -> str:
    """The complete, optimized pipeline for redacting sensitive text."""
    if not text or not isinstance(text, str) or not NLP: 
        return text

    # Process with spaCy once at the beginning for efficiency
    doc = NLP(text)
    redacted_text = text

    # 1. Apply contextual rules first, as they are highly specific
    redacted_text = _apply_contextual_rules(redacted_text, file_name, doc)

    # 2. Re-process with NLP to get accurate entities after initial redactions
    doc = NLP(redacted_text)
    
    # 3. Apply NER (AI-based) redaction
    for ent in reversed(doc.ents):
        if ent.label_ in config.NER_REDACTION_RULES:
            # Check to avoid redacting something already tagged
            if "<" not in ent.text:
                tag = config.NER_REDACTION_RULES[ent.label_]
                audit_log.log_redaction(file_name, ent.text, tag, f"NER: {ent.label_}")
                redacted_text = redacted_text[:ent.start_char] + tag + redacted_text[ent.end_char:]

    # 4. Apply pattern-based (Regex) redaction
    for pattern, tag in config.PII_REDACTION_RULES.items():
        matches = list(pattern.finditer(redacted_text))
        for match in reversed(matches):
            if "<" not in match.group(0): # Avoid double-redacting
                start, end = match.span()
                audit_log.log_redaction(file_name, match.group(0), tag, "Regex")
                redacted_text = redacted_text[:start] + tag + redacted_text[end:]
            
    # 5. Apply keyword-based redaction
    for keyword in config.SENSITIVE_KEYWORDS:
        pattern = re.compile(r'\b' + re.escape(keyword) + r'\b', re.IGNORECASE)
        matches = list(pattern.finditer(redacted_text))
        for match in reversed(matches):
             if "<" not in match.group(0):
                start, end = match.span()
                tag = config.KEYWORD_REPLACEMENT_TAG
                audit_log.log_redaction(file_name, match.group(0), tag, "Keyword")
                redacted_text = redacted_text[:start] + tag + redacted_text[end:]

    return redacted_text

# --- VISUAL DATA CLEANSING (IMAGES) ---

def _find_and_redact_signatures(image):
    """Computer vision function to find and redact handwritten signatures."""
    img_gray = cv2.cvtColor(image, cv2.COLOR_BGR2GRAY)
    _, thresh = cv2.threshold(img_gray, 120, 255, cv2.THRESH_BINARY_INV)
    contours, _ = cv2.findContours(thresh, cv2.RETR_EXTERNAL, cv2.CHAIN_APPROX_SIMPLE)
    for contour in contours:
        area = cv2.contourArea(contour)
        x, y, w, h = cv2.boundingRect(contour)
        aspect_ratio = float(w) / h if h != 0 else 0
        # Tuned parameters for better detection
        if 500 < area < 150000 and 1.2 < aspect_ratio < 15.0:
            # Redact with a solid black box for maximum security
            cv2.rectangle(image, (x, y), (x + w, y + h), (0, 0, 0), -1)
    return image

def _find_and_redact_logos(image, gray_image):
    """Function to find and redact custom logos."""
    for logo in CUSTOM_LOGOS:
        w, h = logo.shape[::-1]
        res = cv2.matchTemplate(gray_image, logo, cv2.TM_CCOEFF_NORMED)
        threshold = 0.8
        loc = np.where(res >= threshold)
        for pt in zip(*loc[::-1]):
            top_left = (int(pt[0]), int(pt[1]))
            bottom_right = (int(pt[0] + w), int(pt[1] + h))
            cv2.rectangle(image, top_left, bottom_right, (0, 0, 0), -1)
    return image

def _load_logos():
    """Loads custom logos from the custom_data directory."""
    logos = []
    if os.path.exists(config.LOGO_DIR):
        for logo_file in os.listdir(config.LOGO_DIR):
            logo_path = os.path.join(config.LOGO_DIR, logo_file)
            logo = cv2.imread(logo_path, 0)
            if logo is not None:
                logos.append(logo)
    return logos
CUSTOM_LOGOS = _load_logos()

# --- FILE-SPECIFIC CLEANSING FUNCTIONS ---

def cleanse_image(input_path: str, output_path: str):
    """The complete cleansing pipeline for image files."""
    original_image = cv2.imread(str(input_path))
    if original_image is None: return

    gray = cv2.cvtColor(original_image, cv2.COLOR_BGR2GRAY)
    
    # 1. Redact visual elements
    original_image = _find_and_redact_logos(original_image, gray)
    original_image = _find_and_redact_signatures(original_image)
    
    face_cascade = cv2.CascadeClassifier(cv2.data.haarcascades + 'haarcascade_frontalface_default.xml')
    faces = face_cascade.detectMultiScale(gray, 1.1, 4)
    for (x, y, w, h) in faces:
        face_roi = original_image[y:y+h, x:x+w]
        original_image[y:y+h, x:x+w] = cv2.GaussianBlur(face_roi, (99, 99), 30) # Blur faces

    # 2. Perform OCR to find and redact sensitive text
    ocr_data = pytesseract.image_to_data(original_image, output_type=pytesseract.Output.DICT)
    
    for i in range(len(ocr_data['text'])):
        word_text = ocr_data['text'][i]
        if int(ocr_data['conf'][i]) > 40 and _is_text_sensitive_advanced(word_text):
            x, y, w, h = ocr_data['left'][i], ocr_data['top'][i], ocr_data['width'][i], ocr_data['height'][i]
            # IMPROVED: Use an irreversible black box instead of a blur for text
            cv2.rectangle(original_image, (x, y), (x + w, y + h), (0, 0, 0), -1)
            
    cv2.imwrite(str(output_path), original_image)

def cleanse_pdf(input_path: str, output_path: str):
    doc = fitz.open(input_path)
    file_name = os.path.basename(input_path)
    for page in doc:
        # Get all text blocks with coordinates
        blocks = page.get_text("dict")["blocks"]
        for b in blocks:
            if b['type'] == 0: # text block
                for l in b["lines"]:
                    for s in l["spans"]:
                        text_to_check = s["text"]
                        redacted_text = redact_text_pipeline(text_to_check, file_name)
                        if text_to_check != redacted_text:
                            # If text is sensitive, add a redaction annotation
                            page.add_redact_annot(s["bbox"], fill=(0,0,0))
        
        # IMPROVED: Securely scrub the text AND the area underneath it.
        page.apply_redactions(images=fitz.PDF_REDACT_IMAGE_PIXELS)

    doc.save(str(output_path), garbage=4, deflate=True, clean=True)
    doc.close()

def cleanse_presentation(input_path: str, output_path: str):
    prs = Presentation(input_path)
    file_name = os.path.basename(input_path)
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    # Redact paragraph by paragraph to maintain structure
                    full_text = "".join(run.text for run in para.runs)
                    redacted_text = redact_text_pipeline(full_text, file_name)
                    if full_text != redacted_text:
                        # Clear existing runs and replace with the single redacted string
                        for i in range(len(para.runs)):
                            para.runs[i].text = ''
                        if para.runs:
                            para.runs[0].text = redacted_text
                        else:
                            para.add_run().text = redacted_text
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        cell.text = redact_text_pipeline(cell.text, file_name)
    prs.save(output_path)


def cleanse_spreadsheet(input_path: str, output_path: str):
    """Cleanses a spreadsheet with a retry mechanism to handle file locks."""
    df = pd.read_excel(input_path) if str(input_path).endswith('.xlsx') else pd.read_csv(input_path)
    file_name = os.path.basename(input_path)
    cleansed_df = df.map(lambda x: redact_text_pipeline(x, file_name) if isinstance(x, str) else x)
    
    # --- NEW: Retry Logic for Saving the File ---
    max_retries = 2
    for attempt in range(max_retries):
        try:
            if str(output_path).endswith('.xlsx'):
                cleansed_df.to_excel(output_path, index=False)
            else:
                cleansed_df.to_csv(output_path, index=False)
            
            # If save is successful, break the loop
            return
            
        except PermissionError as e:
            if attempt < max_retries - 1:
                print(f"   -> WARNING: Could not save {output_path}. File may be open. Retrying in 5 seconds...")
                time.sleep(5) # Wait for 5 seconds
            else:
                # If all retries fail, raise the error to be caught by main.py
                raise e


def cleanse_document(input_path: str, output_path: str):
    doc = Document(input_path)
    file_name = os.path.basename(input_path)
    for para in doc.paragraphs:
        full_text = para.text
        redacted_text = redact_text_pipeline(full_text, file_name)
        if full_text != redacted_text:
            # Clear existing runs and replace content
            for i in range(len(para.runs)):
                para.runs[i].text = ''
            para.text = redacted_text
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                cell.text = redact_text_pipeline(cell.text, file_name)
    doc.save(output_path)

# --- DISPATCH TABLE ---
CLEANSER_DISPATCH = {
    '.png': cleanse_image, '.jpg': cleanse_image, '.jpeg': cleanse_image,
    '.csv': cleanse_spreadsheet, '.xlsx': cleanse_spreadsheet,
    '.pdf': cleanse_pdf, '.pptx': cleanse_presentation, '.docx': cleanse_document,
}
'''
import cv2
import fitz
import pandas as pd
import pytesseract
import re
import os
import time
import numpy as np
import spacy
from pptx import Presentation
from docx import Document
from typing import List, Dict, Tuple, Optional
from . import config, audit_log

# Load AI model with error handling
try:
    NLP = spacy.load("en_core_web_sm")
except OSError:
    print("spaCy model 'en_core_web_sm' not found. Please run 'python -m spacy download en_core_web_sm'")
    NLP = None

pytesseract.pytesseract.tesseract_cmd = config.TESSERACT_CMD_PATH

# --- WHITELIST MANAGEMENT ---

def load_whitelist() -> set:
    """Load whitelist from file if it exists."""
    whitelist = set(config.WHITELIST_KEYWORDS)
    if os.path.exists(config.WHITELIST_FILE):
        with open(config.WHITELIST_FILE, 'r') as f:
            for line in f:
                word = line.strip().lower()
                if word:
                    whitelist.add(word)
    return whitelist

WHITELIST = load_whitelist()

# --- ENHANCED HELPER FUNCTIONS ---

def _calculate_text_entropy(text: str) -> float:
    """Calculate Shannon entropy of text to detect random strings (like keys/tokens)."""
    if not text:
        return 0
    
    prob = [float(text.count(c)) / len(text) for c in set(text)]
    entropy = -sum(p * np.log2(p) for p in prob if p > 0)
    return entropy

def _is_likely_random_string(text: str, threshold: float = 3.5) -> bool:
    """Detect if a string is likely a random key/token based on entropy."""
    if len(text) < 8:  # Too short to be a key
        return False
    
    # Check if it has high entropy (randomness)
    entropy = _calculate_text_entropy(text)
    
    # Check for patterns common in keys/tokens
    has_mixed_case = any(c.isupper() for c in text) and any(c.islower() for c in text)
    has_numbers = any(c.isdigit() for c in text)
    has_special = any(c in '-_+/=' for c in text)
    
    # High entropy with mixed characters suggests a key/token
    if entropy > threshold and (has_mixed_case or has_numbers or has_special):
        return True
    
    return False

def _is_whitelisted(text: str) -> bool:
    """Check if text should not be redacted."""
    text_lower = text.lower().strip()
    
    # Check exact match
    if text_lower in WHITELIST:
        return True
    
    # Check if it's a common English word (length check to avoid short codes)
    if len(text) > 10 and text_lower in WHITELIST:
        return True
    
    # Check if it's purely numeric and short (likely not sensitive)
    if text.isdigit() and len(text) <= 3:
        return True
    
    return False

def _get_context_window(text: str, start: int, end: int, window_size: int = 50) -> str:
    """Get surrounding context for better decision making."""
    context_start = max(0, start - window_size)
    context_end = min(len(text), end + window_size)
    return text[context_start:context_end]

def _is_text_sensitive_advanced(text: str, context: str = "") -> Tuple[bool, str]:
    """
    Advanced check for sensitive information with context awareness.
    Returns (is_sensitive, reason)
    """
    if not text or not isinstance(text, str):
        return False, ""
    
    # Skip if whitelisted
    if _is_whitelisted(text):
        return False, "whitelisted"
    
    # Check if it's a random string (potential key/token)
    if _is_likely_random_string(text):
        return True, "high_entropy_string"
    
    # Check regex patterns with priority
    for pattern, tag in config.PII_REDACTION_RULES.items():
        if pattern.search(text):
            return True, f"regex_pattern_{tag}"
    
    # Check for sensitive keywords with category awareness
    text_lower = text.lower()
    for category, keywords in config.SENSITIVE_KEYWORDS.items():
        for keyword in keywords:
            if keyword.lower() in text_lower:
                # Additional context check for common words
                if category in ["technical", "medical"] and context:
                    # Check if the context suggests it's actually sensitive
                    context_lower = context.lower()
                    if any(indicator in context_lower for indicator in 
                           ["password", "credential", "login", "auth", "patient", "diagnosis"]):
                        return True, f"keyword_{category}_{keyword}"
                else:
                    return True, f"keyword_{category}_{keyword}"
    
    # NLP-based detection if available
    if NLP:
        doc = NLP(text)
        for ent in doc.ents:
            if ent.label_ in config.NER_REDACTION_RULES:
                # Skip common false positives
                if ent.label_ == "CARDINAL" and ent.text.isdigit() and len(ent.text) <= 2:
                    continue
                if ent.label_ == "DATE" and not any(year in ent.text for year in 
                                                   ["1950", "1960", "1970", "1980", "1990", "2000", "2010", "2020"]):
                    continue
                return True, f"ner_{ent.label_}"
    
    return False, ""

def _apply_contextual_rules_enhanced(text: str, file_name: str, nlp_doc) -> str:
    """Enhanced contextual rules to handle numeric and general text targets."""
    redacted_text = text
    
    for rule in config.CONTEXTUAL_RULES:
        trigger_phrase = rule["trigger_phrase"]
        trigger_pattern = re.compile(re.escape(trigger_phrase) + r'\s*:?\s*', re.IGNORECASE)
        
        for match in trigger_pattern.finditer(redacted_text):
            following_start = match.end()
            
            # Extract the value based on the target type
            value_to_redact = None
            
            if rule["TARGET_TYPE"] == "NUM":
                num_pattern = re.compile(r'[\d,]+(?:\.\d{2})?')
                num_match = num_pattern.search(redacted_text, following_start)
                if num_match:
                    value_to_redact = num_match.group(0)
                    actual_start = num_match.start()
                    actual_end = num_match.end()

            # --- NEW LOGIC FOR HANDLING GENERAL TEXT ON THE SAME LINE ---
            elif rule["TARGET_TYPE"] == "TEXT_ON_LINE":
                # Find the end of the line or the entire remaining string
                line_end = redacted_text.find('\n', following_start)
                if line_end == -1:
                    line_end = len(redacted_text)
                
                # Extract the potential value, stripping whitespace
                extracted_value = redacted_text[following_start:line_end].strip()
                if extracted_value:
                    value_to_redact = extracted_value
                    # Find the precise start of the actual text
                    actual_start = redacted_text.find(value_to_redact, following_start)
                    actual_end = actual_start + len(value_to_redact)
            
            # Apply redaction if a valid value was found and not already redacted
            if value_to_redact and "<" not in redacted_text[actual_start:actual_end]:
                tag = rule["redaction_tag"]
                audit_log.log_redaction(file_name, value_to_redact, tag, f"Contextual: '{trigger_phrase}'")
                redacted_text = redacted_text[:actual_start] + tag + redacted_text[actual_end:]
    
    return redacted_text

def _intelligent_merge_redactions(text: str) -> str:
    """Merge adjacent redaction tags for better readability."""
    # Merge same consecutive tags
    patterns = [
        (r'(<\w+>)\s+\1', r'\1'),  # Merge same tags with space
        (r'(<\w+>)(\1)+', r'\1'),   # Merge repeated same tags
    ]
    
    result = text
    for pattern, replacement in patterns:
        result = re.sub(pattern, replacement, result)
    
    return result

def redact_text_pipeline(text: str, file_name: str) -> str:
    """Enhanced complete pipeline for redacting sensitive text."""
    if not text or not isinstance(text, str):
        return text
    
    # Skip very short text
    if len(text.strip()) < 3:
        return text
    
    # Initialize
    redacted_text = text
    redaction_log = []
    
    # Process with spaCy if available
    doc = NLP(text) if NLP else None
    
    # Phase 1: Apply contextual rules first (highest precision)
    if doc:
        redacted_text = _apply_contextual_rules_enhanced(redacted_text, file_name, doc)
    
    # Phase 2: Apply pattern-based redaction with priority order
    for pattern, tag in config.PII_REDACTION_RULES.items():
        matches = list(pattern.finditer(redacted_text))
        
        # Process matches in reverse to maintain indices
        for match in reversed(matches):
            matched_text = match.group(0)
            
            # Skip if already redacted or whitelisted
            if "<" in matched_text or _is_whitelisted(matched_text):
                continue
            
            start, end = match.span()
            context = _get_context_window(redacted_text, start, end)
            
            # Additional validation for certain patterns
            if tag in ["<hostname>", "<email>", "<url>"]:
                # Check if it's actually a sensitive hostname/email/url
                is_sensitive, reason = _is_text_sensitive_advanced(matched_text, context)
                if not is_sensitive:
                    continue
            
            audit_log.log_redaction(file_name, matched_text, tag, "Regex")
            redacted_text = redacted_text[:start] + tag + redacted_text[end:]
            redaction_log.append((matched_text, tag, "regex"))
    
    # Phase 3: Apply NER-based redaction
    if doc:
        # Re-parse after initial redactions
        doc = NLP(redacted_text)
        
        for ent in reversed(doc.ents):
            if ent.label_ in config.NER_REDACTION_RULES:
                tag = config.NER_REDACTION_RULES[ent.label_]
                
                # Skip None tags or already redacted
                if tag is None or "<" in ent.text:
                    continue
                
                # Skip if whitelisted
                if _is_whitelisted(ent.text):
                    continue
                
                # Additional validation for certain entity types
                if ent.label_ == "PERSON":
                    # Check if it's actually a person name (not a common word)
                    if len(ent.text) < 3 or ent.text.lower() in ["the", "and", "for"]:
                        continue
                
                audit_log.log_redaction(file_name, ent.text, tag, f"NER: {ent.label_}")
                redacted_text = redacted_text[:ent.start_char] + tag + redacted_text[ent.end_char:]
                redaction_log.append((ent.text, tag, f"ner_{ent.label_}"))
    
    # Phase 4: Apply keyword-based redaction (lowest priority)
    for category, keywords in config.SENSITIVE_KEYWORDS.items():
        for keyword in keywords:
            # Create word boundary pattern
            pattern = re.compile(r'\b' + re.escape(keyword) + r'\b', re.IGNORECASE)
            matches = list(pattern.finditer(redacted_text))
            
            for match in reversed(matches):
                if "<" not in match.group(0):
                    start, end = match.span()
                    context = _get_context_window(redacted_text, start, end)
                    
                    # Validate based on context
                    if category in ["company_names", "person_names", "banks", "projects"]:
                        # These are likely always sensitive
                        tag = f"<{category[:-1]}>"  # Remove 's' from category name
                        audit_log.log_redaction(file_name, match.group(0), tag, f"Keyword: {category}")
                        redacted_text = redacted_text[:start] + tag + redacted_text[end:]
                        redaction_log.append((match.group(0), tag, f"keyword_{category}"))
    
    # Phase 5: Clean up and merge adjacent redactions
    redacted_text = _intelligent_merge_redactions(redacted_text)
    
    # Log summary if detailed logging is enabled
    if config.ENABLE_DETAILED_LOGGING and redaction_log:
        print(f"  Redacted {len(redaction_log)} items from {file_name}")
        if config.LOG_REDACTION_REASONS:
            for item, tag, reason in redaction_log[:5]:  # Show first 5
                print(f"    - {item[:20]}... -> {tag} ({reason})")
    
    return redacted_text

# --- ENHANCED VISUAL DATA CLEANSING ---

def _find_and_redact_signatures_enhanced(image):
    """Enhanced signature detection using multiple techniques."""
    img_gray = cv2.cvtColor(image, cv2.COLOR_BGR2GRAY)
    
    # Method 1: Threshold-based detection
    _, thresh = cv2.threshold(img_gray, config.SIGNATURE_DETECTION_PARAMS["threshold_value"], 
                             255, cv2.THRESH_BINARY_INV)
    
    # Method 2: Adaptive threshold for better results
    adaptive_thresh = cv2.adaptiveThreshold(img_gray, 255, cv2.ADAPTIVE_THRESH_GAUSSIAN_C,
                                           cv2.THRESH_BINARY_INV, 11, 2)
    
    # Combine both methods
    combined = cv2.bitwise_or(thresh, adaptive_thresh)
    
    # Find contours
    contours, _ = cv2.findContours(combined, cv2.RETR_EXTERNAL, cv2.CHAIN_APPROX_SIMPLE)
    
    params = config.SIGNATURE_DETECTION_PARAMS
    for contour in contours:
        area = cv2.contourArea(contour)
        x, y, w, h = cv2.boundingRect(contour)
        aspect_ratio = float(w) / h if h != 0 else 0
        
        # Enhanced detection criteria
        if (params["min_area"] < area < params["max_area"] and 
            params["min_aspect_ratio"] < aspect_ratio < params["max_aspect_ratio"]):
            
            # Additional check: signature usually has irregular shape
            hull = cv2.convexHull(contour)
            hull_area = cv2.contourArea(hull)
            solidity = float(area) / hull_area if hull_area > 0 else 0
            
            # Signatures typically have lower solidity (more irregular)
            if solidity < 0.95:
                # Redact with solid black box
                cv2.rectangle(image, (x-5, y-5), (x + w + 5, y + h + 5), (0, 0, 0), -1)
    
    return image

def _find_and_redact_text_regions(image, ocr_data, file_name: str):
    """
    Finds and redacts sensitive text regions on a word-by-word basis for accuracy.
    This prevents redacting entire lines and preserves non-sensitive context.
    """
    redacted_image = image.copy()
    num_words = len(ocr_data['text'])

    for i in range(num_words):
        # Only process words with a confidence level above the threshold
        if int(ocr_data['conf'][i]) > config.OCR_CONFIDENCE_THRESHOLD:
            word = ocr_data['text'][i].strip()

            # Skip empty strings or very short, non-sensitive words
            if not word or len(word) < 3:
                continue

            # Check if THIS SPECIFIC word is sensitive using the advanced checker
            is_sensitive, reason = _is_text_sensitive_advanced(word)

            if is_sensitive:
                # Get the precise coordinates (bounding box) of the sensitive word
                (x, y, w, h) = (ocr_data['left'][i], ocr_data['top'][i], ocr_data['width'][i], ocr_data['height'][i])
                
                # Log the specific word being redacted for auditing
                audit_log.log_redaction(file_name, word, f"<visual_{reason}>", "Image OCR")

                # Draw a black rectangle ONLY over the sensitive word
                # The -1 thickness fills the rectangle completely
                cv2.rectangle(redacted_image, (x, y), (x + w, y + h), (0, 0, 0), -1)
    
    return redacted_image



def _find_and_redact_logos(image, gray_image):
    """Function to find and redact custom logos."""
    for logo in CUSTOM_LOGOS:
        w, h = logo.shape[::-1]
        res = cv2.matchTemplate(gray_image, logo, cv2.TM_CCOEFF_NORMED)
        threshold = 0.8
        loc = np.where(res >= threshold)
        for pt in zip(*loc[::-1]):
            top_left = (int(pt[0]), int(pt[1]))
            bottom_right = (int(pt[0] + w), int(pt[1] + h))
            cv2.rectangle(image, top_left, bottom_right, (0, 0, 0), -1)
    return image

def _load_logos():
    """Loads custom logos from the custom_data directory."""
    logos = []
    if os.path.exists(config.LOGO_DIR):
        for logo_file in os.listdir(config.LOGO_DIR):
            logo_path = os.path.join(config.LOGO_DIR, logo_file)
            logo = cv2.imread(logo_path, 0)
            if logo is not None:
                logos.append(logo)
    return logos

CUSTOM_LOGOS = _load_logos()

# --- FILE-SPECIFIC CLEANSING FUNCTIONS ---

def cleanse_image(input_path: str, output_path: str):
    """The complete cleansing pipeline for image files."""
    original_image = cv2.imread(str(input_path))
    if original_image is None: 
        print(f"Warning: Could not read image file {input_path}. Skipping.")
        return

    gray = cv2.cvtColor(original_image, cv2.COLOR_BGR2GRAY)
    file_name = os.path.basename(input_path)
    
    # 1. Redact non-textual visual elements first
    image_with_visuals_redacted = _find_and_redact_logos(original_image, gray)
    image_with_visuals_redacted = _find_and_redact_signatures_enhanced(image_with_visuals_redacted)
    
    # Face detection and blurring
    face_cascade = cv2.CascadeClassifier(cv2.data.haarcascades + 'haarcascade_frontalface_default.xml')
    faces = face_cascade.detectMultiScale(gray, 
                                        scaleFactor=config.FACE_DETECTION_PARAMS["scale_factor"], 
                                        minNeighbors=config.FACE_DETECTION_PARAMS["min_neighbors"])
    
    for (x, y, w, h) in faces:
        face_roi = image_with_visuals_redacted[y:y+h, x:x+w]
        # Apply a heavy Gaussian blur to anonymize the face
        blurred_face = cv2.GaussianBlur(face_roi, 
                                        config.FACE_DETECTION_PARAMS["blur_kernel_size"], 
                                        config.FACE_DETECTION_PARAMS["blur_sigma"])
        image_with_visuals_redacted[y:y+h, x:x+w] = blurred_face

    # 2. Perform OCR to find and redact sensitive text on the visually-cleansed image
    # We use the already processed image to avoid OCR'ing blurred/redacted areas
    ocr_data = pytesseract.image_to_data(image_with_visuals_redacted, output_type=pytesseract.Output.DICT)
    
    # Use the precise, word-by-word text redaction function
    final_image = _find_and_redact_text_regions(image_with_visuals_redacted, ocr_data, file_name)
            
    cv2.imwrite(str(output_path), final_image)

def cleanse_pdf(input_path: str, output_path: str):
    """Cleanse PDF files with enhanced redaction."""
    doc = fitz.open(input_path)
    file_name = os.path.basename(input_path)
    
    for page in doc:
        # Get all text blocks with coordinates
        blocks = page.get_text("dict")["blocks"]
        for b in blocks:
            if b['type'] == 0:  # text block
                for l in b["lines"]:
                    for s in l["spans"]:
                        text_to_check = s["text"]
                        redacted_text = redact_text_pipeline(text_to_check, file_name)
                        if text_to_check != redacted_text:
                            # If text is sensitive, add a redaction annotation
                            page.add_redact_annot(s["bbox"], fill=(0, 0, 0))
        
        # Apply redactions - securely scrub text AND area underneath
        page.apply_redactions(images=fitz.PDF_REDACT_IMAGE_PIXELS)

    doc.save(str(output_path), garbage=4, deflate=True, clean=True)
    doc.close()

def cleanse_presentation(input_path: str, output_path: str):
    """Cleanse PowerPoint presentations."""
    prs = Presentation(input_path)
    file_name = os.path.basename(input_path)
    
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    # Redact paragraph by paragraph to maintain structure
                    full_text = "".join(run.text for run in para.runs)
                    redacted_text = redact_text_pipeline(full_text, file_name)
                    if full_text != redacted_text:
                        # Clear existing runs and replace with the single redacted string
                        for i in range(len(para.runs)):
                            para.runs[i].text = ''
                        if para.runs:
                            para.runs[0].text = redacted_text
                        else:
                            para.add_run().text = redacted_text
            
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        cell.text = redact_text_pipeline(cell.text, file_name)
    
    prs.save(output_path)

def cleanse_spreadsheet(input_path: str, output_path: str):
    """Cleanses a spreadsheet with a retry mechanism to handle file locks."""
    # Read the file based on extension
    if str(input_path).endswith('.xlsx'):
        df = pd.read_excel(input_path)
    else:
        df = pd.read_csv(input_path)
    
    file_name = os.path.basename(input_path)
    
    # Apply redaction to all string cells
    cleansed_df = df.map(lambda x: redact_text_pipeline(str(x), file_name) if pd.notna(x) else x)
    
    # Retry logic for saving the file
    max_retries = 3
    for attempt in range(max_retries):
        try:
            if str(output_path).endswith('.xlsx'):
                cleansed_df.to_excel(output_path, index=False)
            else:
                cleansed_df.to_csv(output_path, index=False)
            
            # If save is successful, break the loop
            return
            
        except PermissionError as e:
            if attempt < max_retries - 1:
                print(f"   -> WARNING: Could not save {output_path}. File may be open. Retrying in 5 seconds...")
                time.sleep(5)  # Wait for 5 seconds
            else:
                # If all retries fail, raise the error to be caught by main.py
                raise e

def cleanse_document(input_path: str, output_path: str):
    """Cleanse Word documents."""
    doc = Document(input_path)
    file_name = os.path.basename(input_path)
    
    # Process paragraphs
    for para in doc.paragraphs:
        full_text = para.text
        redacted_text = redact_text_pipeline(full_text, file_name)
        if full_text != redacted_text:
            # Clear existing runs and replace content
            for i in range(len(para.runs)):
                para.runs[i].text = ''
            if para.runs:
                para.runs[0].text = redacted_text
            else:
                para.text = redacted_text
    
    # Process tables
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                original_text = cell.text
                redacted_text = redact_text_pipeline(original_text, file_name)
                if original_text != redacted_text:
                    cell.text = redacted_text
    
    doc.save(output_path)

# --- DISPATCH TABLE ---
CLEANSER_DISPATCH = {
    '.png': cleanse_image,
    '.jpg': cleanse_image,
    '.jpeg': cleanse_image,
    '.csv': cleanse_spreadsheet,
    '.xlsx': cleanse_spreadsheet,
    '.pdf': cleanse_pdf,
    '.pptx': cleanse_presentation,
    '.docx': cleanse_document,
}